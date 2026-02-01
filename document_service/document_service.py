#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Document Reader Service - Servizio locale per lettura documenti.

Questo servizio legge documenti di vari formati (PDF, Word, Excel, immagini, etc.)
e restituisce il contenuto in formato testo/JSON, utilizzabile da LLM.

Autore: Carlo
Versione: 2.0.0
Porta: 5557
Licenza: MIT

Uso:
    # Avvia il servizio
    python document_service.py

    # Testa con curl
    curl -X POST -F "file=@documento.pdf" http://localhost:5557/read

Formati supportati:
    - Documenti: PDF, Word, Excel, PowerPoint, LibreOffice
    - Immagini: PNG, JPEG, GIF, GIMP, Photoshop, RAW
    - Testo: TXT, Markdown, CSV, JSON, XML, HTML
    - E-book: EPUB, MOBI, AZW
    - Codice: 60+ linguaggi di programmazione
"""

# ============================================================================
# IMPORTS
# ============================================================================
# Librerie standard Python (sempre disponibili)
import os
import sys
import json
import hashlib
import time
import tempfile
import subprocess
import re
import csv
from pathlib import Path
from io import StringIO, BytesIO
from datetime import datetime
from typing import Optional, Dict, Any, List, Callable

# ============================================================================
# VERIFICA DIPENDENZE
# ============================================================================
# Controlliamo quali librerie sono installate per abilitare le funzionalità

# FastAPI - Framework web per le API REST
try:
    from fastapi import FastAPI, File, UploadFile, Form, HTTPException
    from fastapi.middleware.cors import CORSMiddleware
    from fastapi.responses import JSONResponse
    import uvicorn
    HAS_FASTAPI = True
except ImportError:
    HAS_FASTAPI = False

# pypdf - Lettura file PDF
try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    # Prova con il vecchio nome PyPDF2
    try:
        import PyPDF2 as pypdf
        HAS_PYPDF = True
    except ImportError:
        HAS_PYPDF = False

# python-docx - Lettura file Word (.docx)
try:
    from docx import Document as WordDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# openpyxl - Lettura file Excel (.xlsx)
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# python-pptx - Lettura file PowerPoint (.pptx)
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# markdown - Conversione Markdown in HTML
try:
    import markdown
    HAS_MARKDOWN = True
except ImportError:
    HAS_MARKDOWN = False

# beautifulsoup4 - Parsing HTML
try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

# Pillow - Elaborazione immagini
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ebooklib - Lettura EPUB
try:
    from ebooklib import epub
    HAS_EBOOKLIB = True
except ImportError:
    HAS_EBOOKLIB = False


# ============================================================================
# CONFIGURAZIONE
# ============================================================================
# Costanti globali per configurare il comportamento del servizio

# Porta su cui il servizio ascolta le richieste
SERVICE_PORT = 5557

# Cartella dove salvare la cache dei documenti già letti
CACHE_DIR = Path(__file__).parent / ".doc_cache"
CACHE_DIR.mkdir(exist_ok=True)  # Crea la cartella se non esiste

# Dopo quante ore i file in cache vengono considerati "vecchi" e rimossi
CACHE_EXPIRY_HOURS = 24

# Dimensione massima dei file accettati (in megabyte)
MAX_FILE_SIZE_MB = 50


# ============================================================================
# FORMATI SUPPORTATI
# ============================================================================
# Dizionario che mappa ogni estensione al suo tipo di reader
# Struttura: ".estensione": {"name": "Nome visualizzato", "reader": "tipo_reader"}

SUPPORTED_FORMATS: Dict[str, Dict[str, str]] = {
    # -------------------------------------------------------------------------
    # DOCUMENTI OFFICE
    # -------------------------------------------------------------------------
    ".pdf": {"name": "PDF Document", "reader": "pdf"},
    ".docx": {"name": "Word Document", "reader": "docx"},
    ".doc": {"name": "Word Document (Legacy)", "reader": "libreoffice"},
    ".xlsx": {"name": "Excel Spreadsheet", "reader": "xlsx"},
    ".xls": {"name": "Excel Spreadsheet (Legacy)", "reader": "libreoffice"},
    ".pptx": {"name": "PowerPoint Presentation", "reader": "pptx"},
    ".ppt": {"name": "PowerPoint (Legacy)", "reader": "libreoffice"},
    ".rtf": {"name": "Rich Text Format", "reader": "libreoffice"},
    ".wps": {"name": "Microsoft Works", "reader": "libreoffice"},
    ".wpd": {"name": "WordPerfect", "reader": "libreoffice"},

    # -------------------------------------------------------------------------
    # LIBREOFFICE / OPENDOCUMENT
    # -------------------------------------------------------------------------
    ".odt": {"name": "OpenDocument Text", "reader": "libreoffice"},
    ".ods": {"name": "OpenDocument Spreadsheet", "reader": "libreoffice"},
    ".odp": {"name": "OpenDocument Presentation", "reader": "libreoffice"},
    ".odg": {"name": "OpenDocument Graphics", "reader": "libreoffice"},
    ".odf": {"name": "OpenDocument Formula", "reader": "libreoffice"},
    ".odb": {"name": "OpenDocument Database", "reader": "libreoffice"},
    ".ott": {"name": "ODF Text Template", "reader": "libreoffice"},
    ".ots": {"name": "ODF Spreadsheet Template", "reader": "libreoffice"},
    ".otp": {"name": "ODF Presentation Template", "reader": "libreoffice"},
    ".fodt": {"name": "Flat ODF Text", "reader": "libreoffice"},
    ".fods": {"name": "Flat ODF Spreadsheet", "reader": "libreoffice"},
    ".fodp": {"name": "Flat ODF Presentation", "reader": "libreoffice"},

    # -------------------------------------------------------------------------
    # E-BOOK
    # -------------------------------------------------------------------------
    ".epub": {"name": "EPUB E-book", "reader": "epub"},
    ".mobi": {"name": "Kindle E-book", "reader": "ebook"},
    ".azw": {"name": "Amazon Kindle", "reader": "ebook"},
    ".azw3": {"name": "Amazon Kindle KF8", "reader": "ebook"},
    ".fb2": {"name": "FictionBook", "reader": "ebook"},

    # -------------------------------------------------------------------------
    # IMMAGINI - EDITOR GRAFICI
    # -------------------------------------------------------------------------
    ".xcf": {"name": "GIMP Image", "reader": "gimp"},
    ".psd": {"name": "Photoshop Document", "reader": "gimp"},
    ".psb": {"name": "Photoshop Large Document", "reader": "gimp"},
    ".ora": {"name": "OpenRaster", "reader": "gimp"},
    ".kra": {"name": "Krita Document", "reader": "gimp"},
    ".ai": {"name": "Adobe Illustrator", "reader": "gimp"},
    ".eps": {"name": "Encapsulated PostScript", "reader": "gimp"},

    # -------------------------------------------------------------------------
    # IMMAGINI - STANDARD (Pillow)
    # -------------------------------------------------------------------------
    ".png": {"name": "PNG Image", "reader": "image"},
    ".jpg": {"name": "JPEG Image", "reader": "image"},
    ".jpeg": {"name": "JPEG Image", "reader": "image"},
    ".gif": {"name": "GIF Image", "reader": "image"},
    ".bmp": {"name": "Bitmap Image", "reader": "image"},
    ".tiff": {"name": "TIFF Image", "reader": "image"},
    ".tif": {"name": "TIFF Image", "reader": "image"},
    ".webp": {"name": "WebP Image", "reader": "image"},
    ".ico": {"name": "Icon", "reader": "image"},
    ".ppm": {"name": "Portable Pixmap", "reader": "image"},
    ".pgm": {"name": "Portable Graymap", "reader": "image"},
    ".pbm": {"name": "Portable Bitmap", "reader": "image"},
    ".tga": {"name": "Targa", "reader": "image"},

    # -------------------------------------------------------------------------
    # IMMAGINI - RAW (Fotocamere)
    # -------------------------------------------------------------------------
    ".raw": {"name": "RAW Image", "reader": "raw"},
    ".cr2": {"name": "Canon RAW", "reader": "raw"},
    ".cr3": {"name": "Canon RAW 3", "reader": "raw"},
    ".nef": {"name": "Nikon RAW", "reader": "raw"},
    ".arw": {"name": "Sony RAW", "reader": "raw"},
    ".orf": {"name": "Olympus RAW", "reader": "raw"},
    ".rw2": {"name": "Panasonic RAW", "reader": "raw"},
    ".pef": {"name": "Pentax RAW", "reader": "raw"},
    ".dng": {"name": "Digital Negative", "reader": "raw"},
    ".raf": {"name": "Fujifilm RAW", "reader": "raw"},

    # -------------------------------------------------------------------------
    # VETTORIALI
    # -------------------------------------------------------------------------
    ".svg": {"name": "SVG Vector", "reader": "svg"},
    ".svgz": {"name": "SVG Compressed", "reader": "svg"},
    ".wmf": {"name": "Windows Metafile", "reader": "libreoffice"},
    ".emf": {"name": "Enhanced Metafile", "reader": "libreoffice"},

    # -------------------------------------------------------------------------
    # TESTO E MARKUP
    # -------------------------------------------------------------------------
    ".txt": {"name": "Plain Text", "reader": "text"},
    ".md": {"name": "Markdown", "reader": "markdown"},
    ".csv": {"name": "CSV", "reader": "csv"},
    ".tsv": {"name": "Tab-Separated Values", "reader": "csv"},
    ".json": {"name": "JSON", "reader": "json"},
    ".xml": {"name": "XML", "reader": "xml"},
    ".html": {"name": "HTML", "reader": "html"},
    ".htm": {"name": "HTML", "reader": "html"},
    ".xhtml": {"name": "XHTML", "reader": "html"},
    ".tex": {"name": "LaTeX", "reader": "text"},
    ".rst": {"name": "reStructuredText", "reader": "text"},
    ".yaml": {"name": "YAML", "reader": "yaml"},
    ".yml": {"name": "YAML", "reader": "yaml"},
    ".toml": {"name": "TOML", "reader": "text"},
    ".ini": {"name": "INI Config", "reader": "text"},
    ".cfg": {"name": "Config File", "reader": "text"},
    ".conf": {"name": "Config File", "reader": "text"},
    ".log": {"name": "Log File", "reader": "text"},
    ".env": {"name": "Environment File", "reader": "text"},

    # -------------------------------------------------------------------------
    # CODICE SORGENTE (selezione dei più comuni)
    # -------------------------------------------------------------------------
    ".py": {"name": "Python", "reader": "code"},
    ".js": {"name": "JavaScript", "reader": "code"},
    ".ts": {"name": "TypeScript", "reader": "code"},
    ".jsx": {"name": "React JSX", "reader": "code"},
    ".tsx": {"name": "React TSX", "reader": "code"},
    ".java": {"name": "Java", "reader": "code"},
    ".kt": {"name": "Kotlin", "reader": "code"},
    ".c": {"name": "C", "reader": "code"},
    ".cpp": {"name": "C++", "reader": "code"},
    ".h": {"name": "C Header", "reader": "code"},
    ".hpp": {"name": "C++ Header", "reader": "code"},
    ".cs": {"name": "C#", "reader": "code"},
    ".go": {"name": "Go", "reader": "code"},
    ".rs": {"name": "Rust", "reader": "code"},
    ".rb": {"name": "Ruby", "reader": "code"},
    ".php": {"name": "PHP", "reader": "code"},
    ".swift": {"name": "Swift", "reader": "code"},
    ".r": {"name": "R", "reader": "code"},
    ".sql": {"name": "SQL", "reader": "code"},
    ".sh": {"name": "Shell Script", "reader": "code"},
    ".bash": {"name": "Bash Script", "reader": "code"},
    ".ps1": {"name": "PowerShell", "reader": "code"},
    ".bat": {"name": "Batch Script", "reader": "code"},
    ".cmd": {"name": "Windows Command", "reader": "code"},
    ".vue": {"name": "Vue.js", "reader": "code"},
    ".svelte": {"name": "Svelte", "reader": "code"},
    ".css": {"name": "CSS", "reader": "code"},
    ".scss": {"name": "SCSS", "reader": "code"},
    ".sass": {"name": "Sass", "reader": "code"},
    ".less": {"name": "Less", "reader": "code"},
    ".lua": {"name": "Lua", "reader": "code"},
    ".dart": {"name": "Dart", "reader": "code"},
    ".scala": {"name": "Scala", "reader": "code"},
    ".hs": {"name": "Haskell", "reader": "code"},
    ".elm": {"name": "Elm", "reader": "code"},
    ".ex": {"name": "Elixir", "reader": "code"},
    ".erl": {"name": "Erlang", "reader": "code"},
    ".clj": {"name": "Clojure", "reader": "code"},
    ".pl": {"name": "Perl", "reader": "code"},
    ".dockerfile": {"name": "Dockerfile", "reader": "code"},
    ".makefile": {"name": "Makefile", "reader": "code"},
    ".tf": {"name": "Terraform", "reader": "code"},
    ".proto": {"name": "Protocol Buffers", "reader": "code"},
    ".graphql": {"name": "GraphQL", "reader": "code"},
}


# ============================================================================
# CLASSE: DocumentCache
# ============================================================================
class DocumentCache:
    """
    Gestisce la cache dei documenti già letti.

    La cache evita di rileggere lo stesso documento più volte,
    risparmiando tempo e risorse. I documenti sono identificati
    tramite il loro hash MD5.

    Attributi:
        cache_dir: Cartella dove salvare i file di cache
        index_file: File JSON con l'indice della cache
        index: Dizionario con i metadati dei file in cache

    Esempio:
        cache = DocumentCache(Path("./cache"))

        # Controlla se un documento è in cache
        result = cache.get("abc123hash")

        # Salva un risultato in cache
        cache.set("abc123hash", {"text": "contenuto..."})
    """

    def __init__(self, cache_dir: Path) -> None:
        """
        Inizializza la cache.

        Args:
            cache_dir: Percorso della cartella cache
        """
        self.cache_dir = cache_dir
        self.cache_dir.mkdir(exist_ok=True)
        self.index_file = cache_dir / "index.json"
        self.index = self._load_index()

    def _load_index(self) -> Dict[str, Any]:
        """
        Carica l'indice della cache dal file JSON.

        Returns:
            Dizionario con i metadati della cache
        """
        if self.index_file.exists():
            try:
                return json.loads(self.index_file.read_text(encoding="utf-8"))
            except (json.JSONDecodeError, OSError):
                # Se il file è corrotto, ricomincia da zero
                return {}
        return {}

    def _save_index(self) -> None:
        """Salva l'indice della cache su disco."""
        self.index_file.write_text(
            json.dumps(self.index, indent=2, ensure_ascii=False),
            encoding="utf-8"
        )

    def get_hash(self, file_bytes: bytes) -> str:
        """
        Calcola l'hash MD5 di un file.

        L'hash viene usato come identificatore univoco del file.
        Due file identici avranno lo stesso hash.

        Args:
            file_bytes: Contenuto del file in bytes

        Returns:
            Stringa esadecimale con l'hash MD5
        """
        return hashlib.md5(file_bytes).hexdigest()

    def get(self, file_hash: str) -> Optional[Dict[str, Any]]:
        """
        Recupera un documento dalla cache.

        Args:
            file_hash: Hash MD5 del documento

        Returns:
            Dizionario con il risultato, o None se non in cache/scaduto
        """
        if file_hash in self.index:
            entry = self.index[file_hash]
            # Controlla se il documento è ancora valido (non scaduto)
            age_seconds = time.time() - entry.get("timestamp", 0)
            max_age_seconds = CACHE_EXPIRY_HOURS * 3600

            if age_seconds < max_age_seconds:
                return entry.get("result")
        return None

    def set(self, file_hash: str, result: Dict[str, Any]) -> None:
        """
        Salva un documento nella cache.

        Args:
            file_hash: Hash MD5 del documento
            result: Risultato da salvare
        """
        self.index[file_hash] = {
            "timestamp": time.time(),
            "result": result
        }
        self._save_index()

    def cleanup(self) -> int:
        """
        Rimuove i documenti scaduti dalla cache.

        Returns:
            Numero di documenti rimossi
        """
        now = time.time()
        max_age_seconds = CACHE_EXPIRY_HOURS * 3600

        # Trova le entry scadute
        expired = [
            h for h, entry in self.index.items()
            if now - entry.get("timestamp", 0) > max_age_seconds
        ]

        # Rimuovi le entry scadute
        for h in expired:
            del self.index[h]

        self._save_index()
        return len(expired)


# ============================================================================
# CLASSE: DocumentReader
# ============================================================================
class DocumentReader:
    """
    Classe principale per leggere documenti di vari formati.

    Questa classe contiene tutti i metodi per leggere i diversi tipi
    di file. Ogni formato ha il suo metodo dedicato (es. _read_pdf,
    _read_docx, etc.).

    Attributi:
        cache: Istanza di DocumentCache per la cache
        available_readers: Dizionario con i reader disponibili

    Esempio:
        reader = DocumentReader()

        # Leggi un file
        with open("documento.pdf", "rb") as f:
            result = reader.read(f.read(), "documento.pdf")

        print(result["full_text"])
    """

    def __init__(self) -> None:
        """Inizializza il reader e verifica le dipendenze."""
        self.cache = DocumentCache(CACHE_DIR)
        self.available_readers = self._check_readers()

    def _check_readers(self) -> Dict[str, bool]:
        """
        Verifica quali reader sono disponibili sul sistema.

        Controlla sia le librerie Python che i programmi esterni
        (LibreOffice, GIMP, ImageMagick, etc.).

        Returns:
            Dizionario {nome_reader: disponibile}
        """
        import shutil  # Per cercare programmi nel PATH

        readers = {
            # Librerie Python
            "pdf": HAS_PYPDF,
            "docx": HAS_DOCX,
            "xlsx": HAS_OPENPYXL,
            "pptx": HAS_PPTX,
            "markdown": HAS_MARKDOWN,
            "html": HAS_BS4,
            "image": HAS_PIL,
            "epub": HAS_EBOOKLIB,

            # Sempre disponibili (usano solo libreria standard)
            "text": True,
            "csv": True,
            "json": True,
            "xml": True,
            "code": True,
            "yaml": True,
        }

        # Controlla programmi esterni
        # LibreOffice: per formati Office legacy e OpenDocument
        readers["libreoffice"] = (
            shutil.which("soffice") is not None or
            shutil.which("libreoffice") is not None
        )

        # GIMP: per file XCF, PSD, etc.
        readers["gimp"] = shutil.which("gimp") is not None

        # ImageMagick: alternativa per immagini
        readers["imagemagick"] = (
            shutil.which("convert") is not None or
            shutil.which("magick") is not None
        )

        # SVG: cairosvg o inkscape
        try:
            import cairosvg
            readers["svg"] = True
        except ImportError:
            readers["svg"] = (
                shutil.which("inkscape") is not None or
                shutil.which("rsvg-convert") is not None
            )

        # RAW: rawpy o ImageMagick
        try:
            import rawpy
            readers["raw"] = True
        except ImportError:
            readers["raw"] = readers.get("imagemagick", False)

        # Calibre: per e-book diversi da EPUB
        readers["ebook"] = shutil.which("ebook-convert") is not None

        return readers

    def get_supported_formats(self) -> Dict[str, Dict[str, Any]]:
        """
        Restituisce l'elenco dei formati supportati con lo stato.

        Returns:
            Dizionario {estensione: {name, available}}
        """
        result = {}

        for ext, info in SUPPORTED_FORMATS.items():
            reader = info["reader"]

            # Determina se il reader è disponibile
            if reader == "gimp":
                # GIMP o ImageMagick come fallback
                available = (
                    self.available_readers.get("gimp", False) or
                    self.available_readers.get("imagemagick", False)
                )
            elif reader == "raw":
                available = (
                    self.available_readers.get("raw", False) or
                    self.available_readers.get("imagemagick", False)
                )
            elif reader == "ebook":
                available = (
                    self.available_readers.get("ebook", False) or
                    self.available_readers.get("epub", False)
                )
            else:
                available = self.available_readers.get(reader, False)

            result[ext] = {
                "name": info["name"],
                "available": available
            }

        return result

    # -------------------------------------------------------------------------
    # LETTORI SPECIFICI PER FORMATO
    # -------------------------------------------------------------------------

    def _read_pdf(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file PDF ed estrae il testo.

        Args:
            file_bytes: Contenuto del file PDF

        Returns:
            Dizionario con testo, metadati e pagine
        """
        if not HAS_PYPDF:
            return {"error": "pypdf non installato. Installa con: pip install pypdf"}

        try:
            # Apri il PDF dalla memoria
            reader = pypdf.PdfReader(BytesIO(file_bytes))

            # Estrai testo da ogni pagina
            pages_content = []
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text and text.strip():
                    pages_content.append({
                        "page": page_num,
                        "text": text.strip()
                    })

            # Estrai metadati
            metadata = {}
            if reader.metadata:
                metadata_keys = {
                    "/Title": "title",
                    "/Author": "author",
                    "/Subject": "subject",
                    "/Creator": "creator",
                    "/Producer": "producer"
                }
                for pdf_key, clean_key in metadata_keys.items():
                    if pdf_key in reader.metadata:
                        metadata[clean_key] = str(reader.metadata[pdf_key])

            # Combina tutto il testo
            full_text = "\n\n".join(p["text"] for p in pages_content)

            return {
                "format": "PDF",
                "pages": len(reader.pages),
                "metadata": metadata,
                "content": pages_content,
                "full_text": full_text
            }

        except Exception as e:
            return {"error": f"Errore lettura PDF: {str(e)}"}

    def _read_docx(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file Word (.docx) ed estrae testo e tabelle.

        Args:
            file_bytes: Contenuto del file DOCX

        Returns:
            Dizionario con paragrafi, tabelle e metadati
        """
        if not HAS_DOCX:
            return {"error": "python-docx non installato. Installa con: pip install python-docx"}

        try:
            doc = WordDocument(BytesIO(file_bytes))

            # Estrai paragrafi
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        "text": para.text,
                        "style": para.style.name if para.style else "Normal"
                    })

            # Estrai tabelle
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)

            # Estrai metadati
            metadata = {}
            if doc.core_properties:
                props = doc.core_properties
                if props.title:
                    metadata["title"] = props.title
                if props.author:
                    metadata["author"] = props.author
                if props.subject:
                    metadata["subject"] = props.subject
                if props.created:
                    metadata["created"] = str(props.created)

            # Testo completo
            full_text = "\n".join(p["text"] for p in paragraphs)

            return {
                "format": "Word Document",
                "paragraphs_count": len(paragraphs),
                "tables_count": len(tables),
                "metadata": metadata,
                "paragraphs": paragraphs,
                "tables": tables,
                "full_text": full_text
            }

        except Exception as e:
            return {"error": f"Errore lettura DOCX: {str(e)}"}

    def _read_xlsx(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file Excel (.xlsx) ed estrae i dati.

        Args:
            file_bytes: Contenuto del file XLSX

        Returns:
            Dizionario con fogli e dati
        """
        if not HAS_OPENPYXL:
            return {"error": "openpyxl non installato. Installa con: pip install openpyxl"}

        try:
            # data_only=True per ottenere i valori calcolati, non le formule
            wb = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)

            sheets = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]

                # Leggi tutte le righe
                data = []
                for row in ws.iter_rows(values_only=True):
                    # Converti None in stringa vuota
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    # Salta righe completamente vuote
                    if any(cell.strip() for cell in row_data):
                        data.append(row_data)

                sheets[sheet_name] = {
                    "rows": len(data),
                    "cols": len(data[0]) if data else 0,
                    "data": data
                }

            # Converti in testo leggibile
            text_parts = []
            for name, sheet in sheets.items():
                text_parts.append(f"=== Foglio: {name} ===")
                for row in sheet["data"]:
                    text_parts.append("\t".join(row))

            return {
                "format": "Excel Spreadsheet",
                "sheets_count": len(sheets),
                "sheets": sheets,
                "full_text": "\n".join(text_parts)
            }

        except Exception as e:
            return {"error": f"Errore lettura XLSX: {str(e)}"}

    def _read_pptx(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file PowerPoint (.pptx) ed estrae testo e note.

        Args:
            file_bytes: Contenuto del file PPTX

        Returns:
            Dizionario con slide e testo
        """
        if not HAS_PPTX:
            return {"error": "python-pptx non installato. Installa con: pip install python-pptx"}

        try:
            prs = Presentation(BytesIO(file_bytes))

            slides = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_content = {
                    "number": i,
                    "texts": [],
                    "notes": ""
                }

                # Estrai testo dalle forme
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["texts"].append(shape.text.strip())

                # Estrai note
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_content["notes"] = slide.notes_slide.notes_text_frame.text

                slides.append(slide_content)

            # Testo completo
            text_parts = []
            for slide in slides:
                text_parts.append(f"=== Slide {slide['number']} ===")
                text_parts.extend(slide["texts"])
                if slide["notes"]:
                    text_parts.append(f"[Note: {slide['notes']}]")

            return {
                "format": "PowerPoint Presentation",
                "slides_count": len(slides),
                "slides": slides,
                "full_text": "\n".join(text_parts)
            }

        except Exception as e:
            return {"error": f"Errore lettura PPTX: {str(e)}"}

    def _read_text(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file di testo semplice.

        Prova diverse codifiche per gestire file con encoding diversi.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con testo e info
        """
        # Prova diverse codifiche in ordine di probabilità
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]

        for encoding in encodings:
            try:
                text = file_bytes.decode(encoding)
                return {
                    "format": "Plain Text",
                    "encoding": encoding,
                    "lines": text.count("\n") + 1,
                    "characters": len(text),
                    "full_text": text
                }
            except UnicodeDecodeError:
                continue

        return {"error": "Impossibile decodificare il file con le codifiche supportate"}

    def _read_markdown(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file Markdown ed estrae la struttura.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con testo, headers e HTML
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        result["format"] = "Markdown"
        text = result["full_text"]

        # Estrai gli header (# Header)
        headers = re.findall(r'^(#{1,6})\s+(.+)$', text, re.MULTILINE)
        result["headers"] = [
            {"level": len(h[0]), "text": h[1]}
            for h in headers
        ]

        # Converti in HTML se possibile
        if HAS_MARKDOWN:
            result["html"] = markdown.markdown(text)

        return result

    def _read_csv(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file CSV/TSV e rileva automaticamente il delimitatore.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con dati e struttura
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        try:
            text = result["full_text"]

            # Rileva il delimitatore più comune
            sample = text[:2048]  # Usa i primi 2KB per rilevare
            delimiters = [",", ";", "\t", "|"]
            delimiter = max(delimiters, key=lambda d: sample.count(d))

            # Parsa il CSV
            reader = csv.reader(StringIO(text), delimiter=delimiter)
            rows = list(reader)

            return {
                "format": "CSV",
                "delimiter": repr(delimiter),  # Mostra il carattere in modo leggibile
                "rows": len(rows),
                "columns": len(rows[0]) if rows else 0,
                "headers": rows[0] if rows else [],
                "data": rows[1:] if len(rows) > 1 else [],
                "full_text": text
            }

        except Exception as e:
            return {"error": f"Errore lettura CSV: {str(e)}"}

    def _read_json(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file JSON e ne parsa la struttura.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con dati JSON parsati
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        try:
            data = json.loads(result["full_text"])

            return {
                "format": "JSON",
                "type": type(data).__name__,  # list, dict, etc.
                "keys": list(data.keys()) if isinstance(data, dict) else None,
                "length": len(data) if isinstance(data, (list, dict)) else None,
                "data": data,
                "full_text": json.dumps(data, indent=2, ensure_ascii=False)
            }

        except json.JSONDecodeError as e:
            return {
                "error": f"JSON non valido: {str(e)}",
                "full_text": result["full_text"]
            }

    def _read_xml(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file XML.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con info sulla struttura XML
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        result["format"] = "XML"
        text = result["full_text"]

        # Trova l'elemento root
        root_match = re.search(r'<(\w+)[^>]*>', text)
        if root_match:
            result["root_element"] = root_match.group(1)

        return result

    def _read_html(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file HTML ed estrae il testo.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con testo estratto e struttura
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        result["format"] = "HTML"

        if HAS_BS4:
            soup = BeautifulSoup(result["full_text"], "html.parser")

            # Estrai titolo
            title = soup.find("title")
            if title:
                result["title"] = title.get_text()

            # Estrai solo il testo leggibile
            result["plain_text"] = soup.get_text(separator="\n", strip=True)

            # Estrai gli header
            headers = []
            for level in range(1, 7):
                for h in soup.find_all(f"h{level}"):
                    headers.append({"level": level, "text": h.get_text()})
            result["headers"] = headers

        return result

    def _read_yaml(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file YAML.

        Args:
            file_bytes: Contenuto del file

        Returns:
            Dizionario con dati YAML parsati
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        result["format"] = "YAML"

        try:
            import yaml
            data = yaml.safe_load(result["full_text"])
            result["data"] = data
            result["type"] = type(data).__name__
            if isinstance(data, dict):
                result["keys"] = list(data.keys())
        except ImportError:
            pass  # PyYAML non installato, OK
        except Exception:
            pass  # YAML malformato, restituisci comunque il testo

        return result

    def _read_code(self, file_bytes: bytes, language: str = "unknown") -> Dict[str, Any]:
        """
        Legge un file di codice sorgente.

        Args:
            file_bytes: Contenuto del file
            language: Linguaggio di programmazione

        Returns:
            Dizionario con codice e statistiche base
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        result["format"] = f"Source Code ({language})"
        result["language"] = language
        text = result["full_text"]

        # Statistiche base per alcuni linguaggi
        if language in ("py", "python"):
            result["functions"] = len(re.findall(r'^def\s+\w+', text, re.MULTILINE))
            result["classes"] = len(re.findall(r'^class\s+\w+', text, re.MULTILINE))
            result["imports"] = len(re.findall(r'^(?:import|from)\s+', text, re.MULTILINE))
        elif language in ("js", "javascript", "ts", "typescript"):
            result["functions"] = len(re.findall(r'function\s+\w+|const\s+\w+\s*=\s*(?:async\s*)?\(', text))
            result["classes"] = len(re.findall(r'class\s+\w+', text))
        elif language in ("java", "cs", "csharp"):
            result["classes"] = len(re.findall(r'class\s+\w+', text))
            result["methods"] = len(re.findall(r'(?:public|private|protected)\s+\w+\s+\w+\s*\(', text))

        return result

    def _read_image(self, file_bytes: bytes, ext: str = "") -> Dict[str, Any]:
        """
        Legge i metadati di un'immagine.

        Args:
            file_bytes: Contenuto del file immagine
            ext: Estensione del file

        Returns:
            Dizionario con metadati immagine
        """
        if not HAS_PIL:
            return {"error": "Pillow non installato. Installa con: pip install Pillow"}

        try:
            img = Image.open(BytesIO(file_bytes))

            result = {
                "format": f"Image ({img.format or ext.upper().replace('.', '')})",
                "width": img.size[0],
                "height": img.size[1],
                "mode": img.mode,  # RGB, RGBA, L (grayscale), etc.
                "size_bytes": len(file_bytes),
                "size_kb": round(len(file_bytes) / 1024, 2),
            }

            # Cerca dati EXIF (metadati fotografici)
            if hasattr(img, "_getexif") and callable(img._getexif):
                try:
                    exif = img._getexif()
                    if exif:
                        exif_tags = {
                            271: "camera_make",
                            272: "camera_model",
                            306: "datetime",
                            37386: "focal_length",
                        }
                        result["exif"] = {}
                        for tag_id, name in exif_tags.items():
                            if tag_id in exif:
                                result["exif"][name] = str(exif[tag_id])[:100]
                except Exception:
                    pass

            # Descrizione testuale
            result["full_text"] = (
                f"Immagine {img.format or ext.upper()}\n"
                f"Dimensioni: {img.size[0]}x{img.size[1]} pixel\n"
                f"Modalità colore: {img.mode}\n"
                f"Dimensione file: {result['size_kb']} KB"
            )

            return result

        except Exception as e:
            return {"error": f"Errore lettura immagine: {str(e)}"}

    def _read_with_libreoffice(self, file_bytes: bytes, ext: str) -> Dict[str, Any]:
        """
        Legge un documento usando LibreOffice in modalità headless.

        Questo metodo è usato per formati che non hanno una libreria
        Python dedicata, come .doc, .xls, .odt, etc.

        Args:
            file_bytes: Contenuto del file
            ext: Estensione del file

        Returns:
            Dizionario con testo estratto
        """
        import shutil

        # Trova l'eseguibile di LibreOffice
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return {
                "error": "LibreOffice non installato.",
                "suggestion": "Installa con: sudo apt install libreoffice (Linux) o winget install TheDocumentFoundation.LibreOffice (Windows)"
            }

        try:
            # Salva il file temporaneamente
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            # Crea cartella per l'output
            output_dir = tempfile.mkdtemp()

            # Converti in testo usando LibreOffice
            subprocess.run(
                [
                    soffice,
                    "--headless",  # Senza interfaccia grafica
                    "--convert-to", "txt:Text",  # Formato output
                    "--outdir", output_dir,
                    tmp_path
                ],
                capture_output=True,
                timeout=60  # Timeout di 60 secondi
            )

            # Leggi il file di output
            txt_file = Path(output_dir) / (Path(tmp_path).stem + ".txt")
            if txt_file.exists():
                text = txt_file.read_text(encoding="utf-8", errors="ignore")

                # Pulizia file temporanei
                os.unlink(tmp_path)
                txt_file.unlink()
                os.rmdir(output_dir)

                return {
                    "format": SUPPORTED_FORMATS.get(ext, {}).get("name", "Document"),
                    "conversion_method": "LibreOffice",
                    "full_text": text
                }
            else:
                return {"error": "Conversione LibreOffice fallita - file output non creato"}

        except subprocess.TimeoutExpired:
            return {"error": "Timeout durante la conversione con LibreOffice"}
        except Exception as e:
            return {"error": f"Errore LibreOffice: {str(e)}"}

    def _read_gimp(self, file_bytes: bytes, ext: str) -> Dict[str, Any]:
        """
        Legge file di editor grafici (XCF, PSD, etc.) usando ImageMagick o GIMP.

        Args:
            file_bytes: Contenuto del file
            ext: Estensione del file

        Returns:
            Dizionario con metadati immagine
        """
        import shutil

        # Prova prima con ImageMagick (più comune)
        magick = shutil.which("magick") or shutil.which("convert")
        if magick:
            return self._read_with_imagemagick(file_bytes, ext)

        # Prova con GIMP
        if shutil.which("gimp"):
            return self._read_with_gimp_cli(file_bytes, ext)

        # Fallback: Pillow per PSD (supporto limitato)
        if ext.lower() == ".psd" and HAS_PIL:
            return self._read_image(file_bytes, ext)

        return {
            "error": f"Nessun programma disponibile per leggere {ext}",
            "suggestions": [
                "sudo apt install imagemagick  # Linux",
                "sudo apt install gimp  # Linux",
                "winget install ImageMagick.ImageMagick  # Windows",
                "winget install GIMP.GIMP  # Windows"
            ]
        }

    def _read_with_imagemagick(self, file_bytes: bytes, ext: str) -> Dict[str, Any]:
        """
        Estrae informazioni da un'immagine usando ImageMagick.

        Args:
            file_bytes: Contenuto del file
            ext: Estensione del file

        Returns:
            Dizionario con metadati
        """
        import shutil

        try:
            # Salva temporaneamente
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            # Usa identify per ottenere info
            magick = shutil.which("magick")
            if magick:
                cmd = ["magick", "identify", "-verbose", tmp_path]
            else:
                cmd = ["identify", "-verbose", tmp_path]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            info_text = result.stdout

            # Parse delle informazioni base
            info = {
                "format": ext.upper().replace(".", ""),
                "conversion_method": "ImageMagick",
            }

            # Estrai dimensioni
            geom_match = re.search(r'Geometry:\s*(\d+)x(\d+)', info_text)
            if geom_match:
                info["width"] = int(geom_match.group(1))
                info["height"] = int(geom_match.group(2))

            # Estrai tipo colore
            type_match = re.search(r'Type:\s*(\w+)', info_text)
            if type_match:
                info["color_type"] = type_match.group(1)

            info["full_text"] = info_text[:5000]  # Limita l'output

            os.unlink(tmp_path)
            return info

        except subprocess.TimeoutExpired:
            return {"error": "Timeout ImageMagick"}
        except Exception as e:
            return {"error": f"Errore ImageMagick: {str(e)}"}

    def _read_with_gimp_cli(self, file_bytes: bytes, ext: str) -> Dict[str, Any]:
        """
        Legge file usando GIMP in modalità batch.

        Args:
            file_bytes: Contenuto del file
            ext: Estensione del file

        Returns:
            Dizionario con info dal file GIMP
        """
        try:
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            # GIMP script per ottenere info base
            # Nota: GIMP batch scripting è complesso, restituiamo info minime
            result = subprocess.run(
                ["gimp", "-i", "-b", f'(gimp-quit 0)'],
                capture_output=True,
                text=True,
                timeout=30
            )

            os.unlink(tmp_path)

            return {
                "format": SUPPORTED_FORMATS.get(ext, {}).get("name", "GIMP Image"),
                "conversion_method": "GIMP",
                "full_text": f"File {ext} - usa GIMP per visualizzare i dettagli"
            }

        except Exception as e:
            return {"error": f"Errore GIMP: {str(e)}"}

    def _read_raw(self, file_bytes: bytes, ext: str) -> Dict[str, Any]:
        """
        Legge file RAW da fotocamere digitali.

        Args:
            file_bytes: Contenuto del file
            ext: Estensione del file

        Returns:
            Dizionario con metadati RAW
        """
        # Prova con rawpy
        try:
            import rawpy

            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            with rawpy.imread(tmp_path) as raw:
                result = {
                    "format": f"RAW Image ({ext.upper().replace('.', '')})",
                    "width": raw.sizes.width,
                    "height": raw.sizes.height,
                }

                # Info fotocamera se disponibili
                if hasattr(raw, "camera_make"):
                    result["camera_make"] = raw.camera_make.decode() if isinstance(raw.camera_make, bytes) else str(raw.camera_make)
                if hasattr(raw, "camera_model"):
                    result["camera_model"] = raw.camera_model.decode() if isinstance(raw.camera_model, bytes) else str(raw.camera_model)

                result["full_text"] = f"RAW {ext.upper()}: {raw.sizes.width}x{raw.sizes.height}"

            os.unlink(tmp_path)
            return result

        except ImportError:
            pass
        except Exception:
            pass

        # Fallback a ImageMagick
        return self._read_with_imagemagick(file_bytes, ext)

    def _read_svg(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file SVG ed estrae informazioni sulla struttura.

        Args:
            file_bytes: Contenuto del file SVG

        Returns:
            Dizionario con info sul documento SVG
        """
        result = self._read_text(file_bytes)
        if "error" in result:
            return result

        result["format"] = "SVG Vector"
        text = result["full_text"]

        # Estrai dimensioni
        width = re.search(r'width=["\']?(\d+)', text)
        height = re.search(r'height=["\']?(\d+)', text)
        viewbox = re.search(r'viewBox=["\']([^"\']+)["\']', text)

        if width:
            result["width"] = width.group(1)
        if height:
            result["height"] = height.group(1)
        if viewbox:
            result["viewbox"] = viewbox.group(1)

        # Conta elementi comuni
        result["elements"] = {
            "paths": len(re.findall(r'<path', text)),
            "circles": len(re.findall(r'<circle', text)),
            "rects": len(re.findall(r'<rect', text)),
            "texts": len(re.findall(r'<text', text)),
            "groups": len(re.findall(r'<g[ >]', text)),
        }

        return result

    def _read_epub(self, file_bytes: bytes) -> Dict[str, Any]:
        """
        Legge un file EPUB ed estrae testo e metadati.

        Args:
            file_bytes: Contenuto del file EPUB

        Returns:
            Dizionario con capitoli e metadati
        """
        if not HAS_EBOOKLIB:
            return {"error": "ebooklib non installato. Installa con: pip install ebooklib"}

        try:
            book = epub.read_epub(BytesIO(file_bytes))

            # Metadati
            title = book.get_metadata('DC', 'title')
            author = book.get_metadata('DC', 'creator')
            language = book.get_metadata('DC', 'language')

            result = {
                "format": "EPUB E-book",
                "title": title[0][0] if title else None,
                "author": author[0][0] if author else None,
                "language": language[0][0] if language else None,
            }

            # Estrai testo dai capitoli
            chapters = []
            text_parts = []

            for item in book.get_items():
                if item.get_type() == 9:  # ITEM_DOCUMENT
                    content = item.get_content().decode('utf-8', errors='ignore')

                    # Rimuovi tag HTML
                    if HAS_BS4:
                        soup = BeautifulSoup(content, 'html.parser')
                        text = soup.get_text(separator='\n', strip=True)
                    else:
                        text = re.sub(r'<[^>]+>', '', content)

                    if text.strip():
                        chapters.append({
                            "name": item.get_name(),
                            "preview": text[:300]  # Anteprima
                        })
                        text_parts.append(text)

            result["chapters_count"] = len(chapters)
            result["chapters"] = chapters[:20]  # Max 20 capitoli nell'output
            result["full_text"] = "\n\n".join(text_parts)[:100000]  # Max 100k caratteri

            return result

        except Exception as e:
            return {"error": f"Errore lettura EPUB: {str(e)}"}

    def _read_ebook(self, file_bytes: bytes, ext: str) -> Dict[str, Any]:
        """
        Legge altri formati e-book usando Calibre.

        Args:
            file_bytes: Contenuto del file
            ext: Estensione del file

        Returns:
            Dizionario con testo estratto
        """
        import shutil

        ebook_convert = shutil.which("ebook-convert")
        if not ebook_convert:
            # Prova EPUB se possibile
            if ext.lower() == ".epub":
                return self._read_epub(file_bytes)

            return {
                "error": "Calibre non installato.",
                "suggestion": "Installa con: sudo apt install calibre (Linux) o scarica da https://calibre-ebook.com/"
            }

        try:
            # Salva temporaneamente
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            # Converti in testo
            txt_path = tmp_path.rsplit(".", 1)[0] + ".txt"
            subprocess.run(
                [ebook_convert, tmp_path, txt_path],
                capture_output=True,
                timeout=120
            )

            if os.path.exists(txt_path):
                text = Path(txt_path).read_text(encoding='utf-8', errors='ignore')

                # Pulizia
                os.unlink(tmp_path)
                os.unlink(txt_path)

                return {
                    "format": SUPPORTED_FORMATS.get(ext, {}).get("name", "E-book"),
                    "conversion_method": "Calibre",
                    "full_text": text
                }

            os.unlink(tmp_path)
            return {"error": "Conversione Calibre fallita"}

        except subprocess.TimeoutExpired:
            return {"error": "Timeout Calibre"}
        except Exception as e:
            return {"error": f"Errore Calibre: {str(e)}"}

    # -------------------------------------------------------------------------
    # METODO PRINCIPALE
    # -------------------------------------------------------------------------

    def read(
        self,
        file_bytes: bytes,
        filename: str,
        use_cache: bool = True
    ) -> Dict[str, Any]:
        """
        Legge un documento e restituisce il contenuto strutturato.

        Questo è il metodo principale da chiamare. Determina automaticamente
        il tipo di file dall'estensione e usa il reader appropriato.

        Args:
            file_bytes: Contenuto del file in bytes
            filename: Nome del file (usato per determinare il formato)
            use_cache: Se True, usa la cache per evitare riletture

        Returns:
            Dizionario con:
                - format: Tipo di documento
                - full_text: Testo estratto
                - Altri campi specifici per formato
                - error: Messaggio di errore se fallisce

        Esempio:
            with open("documento.pdf", "rb") as f:
                result = reader.read(f.read(), "documento.pdf")

            if "error" not in result:
                print(result["full_text"])
        """
        # Controlla la cache
        file_hash = self.cache.get_hash(file_bytes)

        if use_cache:
            cached = self.cache.get(file_hash)
            if cached:
                cached["from_cache"] = True
                return cached

        # Determina il formato dal nome file
        ext = Path(filename).suffix.lower()

        if ext not in SUPPORTED_FORMATS:
            return {
                "error": f"Formato non supportato: {ext}",
                "supported_formats": list(SUPPORTED_FORMATS.keys())
            }

        format_info = SUPPORTED_FORMATS[ext]
        reader_type = format_info["reader"]

        # Mappa dei reader disponibili
        readers: Dict[str, Callable[[], Dict[str, Any]]] = {
            "pdf": lambda: self._read_pdf(file_bytes),
            "docx": lambda: self._read_docx(file_bytes),
            "xlsx": lambda: self._read_xlsx(file_bytes),
            "pptx": lambda: self._read_pptx(file_bytes),
            "text": lambda: self._read_text(file_bytes),
            "markdown": lambda: self._read_markdown(file_bytes),
            "csv": lambda: self._read_csv(file_bytes),
            "json": lambda: self._read_json(file_bytes),
            "xml": lambda: self._read_xml(file_bytes),
            "html": lambda: self._read_html(file_bytes),
            "yaml": lambda: self._read_yaml(file_bytes),
            "code": lambda: self._read_code(file_bytes, ext.lstrip(".")),
            "libreoffice": lambda: self._read_with_libreoffice(file_bytes, ext),
            "image": lambda: self._read_image(file_bytes, ext),
            "gimp": lambda: self._read_gimp(file_bytes, ext),
            "raw": lambda: self._read_raw(file_bytes, ext),
            "svg": lambda: self._read_svg(file_bytes),
            "epub": lambda: self._read_epub(file_bytes),
            "ebook": lambda: self._read_ebook(file_bytes, ext),
        }

        # Esegui il reader appropriato
        if reader_type in readers:
            result = readers[reader_type]()
        else:
            result = {"error": f"Reader non implementato: {reader_type}"}

        # Aggiungi metadati comuni
        result["filename"] = filename
        result["extension"] = ext
        result["size_bytes"] = len(file_bytes)
        result["size_kb"] = round(len(file_bytes) / 1024, 2)
        result["hash"] = file_hash
        result["timestamp"] = datetime.now().isoformat()
        result["from_cache"] = False

        # Salva in cache se non c'è errore
        if use_cache and "error" not in result:
            self.cache.set(file_hash, result)

        return result

    def get_summary(
        self,
        file_bytes: bytes,
        filename: str,
        max_chars: int = 2000
    ) -> str:
        """
        Restituisce un riassunto breve del documento.

        Args:
            file_bytes: Contenuto del file
            filename: Nome del file
            max_chars: Massimo numero di caratteri

        Returns:
            Stringa con riassunto del documento
        """
        result = self.read(file_bytes, filename)

        if "error" in result:
            return f"Errore: {result['error']}"

        full_text = result.get("full_text", "")
        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n...[troncato]"

        return f"[{result.get('format', 'Documento')}] {result.get('filename', '')}\n\n{full_text}"


# ============================================================================
# API SERVICE - FastAPI
# ============================================================================

def create_app() -> "FastAPI":
    """
    Crea e configura l'applicazione FastAPI.

    Returns:
        Istanza FastAPI configurata con tutti gli endpoint
    """
    app = FastAPI(
        title="Document Reader Service",
        description="Servizio locale per lettura documenti di vari formati",
        version="2.0.0",
        docs_url="/docs",  # Swagger UI disponibile su /docs
        redoc_url="/redoc"  # ReDoc disponibile su /redoc
    )

    # Configura CORS per permettere richieste da Open WebUI
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],  # Permetti tutte le origini
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    # Inizializza il reader
    reader = DocumentReader()

    # -------------------------------------------------------------------------
    # ENDPOINT: Home / Health Check
    # -------------------------------------------------------------------------
    @app.get("/", tags=["Info"])
    async def root() -> Dict[str, Any]:
        """
        Health check e informazioni sul servizio.

        Restituisce lo stato del servizio e i formati disponibili.
        """
        formats = reader.get_supported_formats()
        available = [ext for ext, info in formats.items() if info["available"]]

        return {
            "service": "Document Reader Service",
            "version": "2.0.0",
            "status": "running",
            "port": SERVICE_PORT,
            "formats_available": len(available),
            "formats_total": len(formats),
            "documentation": f"http://localhost:{SERVICE_PORT}/docs",
            "endpoints": [
                "POST /read - Legge un documento",
                "POST /extract-text - Estrae solo il testo",
                "POST /get-metadata - Restituisce solo metadati",
                "POST /summary - Riassunto breve",
                "GET /formats - Lista formati supportati",
                "DELETE /cache - Pulisce la cache"
            ]
        }

    # -------------------------------------------------------------------------
    # ENDPOINT: Lista Formati
    # -------------------------------------------------------------------------
    @app.get("/formats", tags=["Info"])
    async def list_formats() -> Dict[str, Dict[str, Any]]:
        """
        Lista tutti i formati supportati con lo stato di disponibilità.

        Ogni formato indica se le dipendenze necessarie sono installate.
        """
        return reader.get_supported_formats()

    # -------------------------------------------------------------------------
    # ENDPOINT: Lettura Documento
    # -------------------------------------------------------------------------
    @app.post("/read", tags=["Documenti"])
    async def read_document(
        file: UploadFile = File(..., description="File documento da leggere"),
        use_cache: bool = Form(default=True, description="Usa cache per risultati")
    ) -> JSONResponse:
        """
        Legge un documento e restituisce il contenuto strutturato.

        Supporta PDF, Word, Excel, PowerPoint, immagini, e-book e molto altro.
        Vedi /formats per l'elenco completo.
        """
        try:
            # Leggi il contenuto del file
            contents = await file.read()

            # Verifica dimensione
            size_mb = len(contents) / (1024 * 1024)
            if size_mb > MAX_FILE_SIZE_MB:
                raise HTTPException(
                    400,
                    f"File troppo grande ({size_mb:.1f}MB). Massimo: {MAX_FILE_SIZE_MB}MB"
                )

            # Leggi il documento
            result = reader.read(
                contents,
                file.filename or "document",
                use_cache=use_cache
            )

            return JSONResponse(result)

        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"Errore lettura: {str(e)}")

    # -------------------------------------------------------------------------
    # ENDPOINT: Estrazione Testo
    # -------------------------------------------------------------------------
    @app.post("/extract-text", tags=["Documenti"])
    async def extract_text(
        file: UploadFile = File(..., description="File da cui estrarre testo")
    ) -> Dict[str, Any]:
        """
        Estrae solo il testo dal documento (senza metadati).

        Utile quando serve solo il contenuto testuale.
        """
        try:
            contents = await file.read()
            result = reader.read(contents, file.filename or "document")

            if "error" in result:
                return {"error": result["error"]}

            return {
                "filename": result.get("filename"),
                "format": result.get("format"),
                "text": result.get("full_text", ""),
                "characters": len(result.get("full_text", ""))
            }

        except Exception as e:
            raise HTTPException(500, f"Errore: {str(e)}")

    # -------------------------------------------------------------------------
    # ENDPOINT: Metadati
    # -------------------------------------------------------------------------
    @app.post("/get-metadata", tags=["Documenti"])
    async def get_metadata(
        file: UploadFile = File(..., description="File di cui estrarre metadati")
    ) -> JSONResponse:
        """
        Restituisce solo i metadati del documento (senza contenuto testuale).

        Utile per ottenere informazioni rapide su un file.
        """
        try:
            contents = await file.read()
            result = reader.read(contents, file.filename or "document")

            # Rimuovi il testo per avere solo metadati
            fields_to_remove = ["full_text", "content", "data", "paragraphs", "slides", "sheets"]
            metadata = {k: v for k, v in result.items() if k not in fields_to_remove}

            return JSONResponse(metadata)

        except Exception as e:
            raise HTTPException(500, f"Errore: {str(e)}")

    # -------------------------------------------------------------------------
    # ENDPOINT: Riassunto
    # -------------------------------------------------------------------------
    @app.post("/summary", tags=["Documenti"])
    async def get_summary(
        file: UploadFile = File(..., description="File da riassumere"),
        max_chars: int = Form(default=2000, description="Massimo caratteri")
    ) -> Dict[str, str]:
        """
        Restituisce un riassunto breve del documento.

        Utile per avere una panoramica veloce del contenuto.
        """
        try:
            contents = await file.read()
            summary = reader.get_summary(contents, file.filename or "document", max_chars)

            return {"summary": summary}

        except Exception as e:
            raise HTTPException(500, f"Errore: {str(e)}")

    # -------------------------------------------------------------------------
    # ENDPOINT: Batch
    # -------------------------------------------------------------------------
    @app.post("/batch", tags=["Documenti"])
    async def batch_read(
        files: List[UploadFile] = File(..., description="File da leggere")
    ) -> Dict[str, List[Dict[str, Any]]]:
        """
        Legge multiple documenti in un'unica richiesta.

        Utile per processare molti file insieme.
        """
        results = []

        for file in files:
            try:
                contents = await file.read()
                result = reader.read(contents, file.filename or "document")
                results.append(result)
            except Exception as e:
                results.append({
                    "filename": file.filename,
                    "error": str(e)
                })

        return {"results": results}

    # -------------------------------------------------------------------------
    # ENDPOINT: Cache
    # -------------------------------------------------------------------------
    @app.delete("/cache", tags=["Sistema"])
    async def clear_cache() -> Dict[str, str]:
        """
        Pulisce la cache dei documenti.

        Rimuove i file più vecchi di 24 ore.
        """
        removed = reader.cache.cleanup()
        return {"message": f"Cache pulita. Rimossi {removed} documenti scaduti."}

    return app


# ============================================================================
# FUNZIONE: Verifica Dipendenze
# ============================================================================

def check_dependencies() -> bool:
    """
    Verifica e mostra lo stato delle dipendenze.

    Returns:
        True se le dipendenze minime sono soddisfatte
    """
    print("\n" + "=" * 60)
    print("VERIFICA DIPENDENZE")
    print("=" * 60 + "\n")

    # Lista delle dipendenze con nome, stato e comando installazione
    deps = [
        ("FastAPI + Uvicorn", HAS_FASTAPI, "pip install fastapi uvicorn python-multipart"),
        ("pypdf (PDF)", HAS_PYPDF, "pip install pypdf"),
        ("python-docx (Word)", HAS_DOCX, "pip install python-docx"),
        ("openpyxl (Excel)", HAS_OPENPYXL, "pip install openpyxl"),
        ("python-pptx (PowerPoint)", HAS_PPTX, "pip install python-pptx"),
        ("Pillow (Immagini)", HAS_PIL, "pip install Pillow"),
        ("markdown", HAS_MARKDOWN, "pip install markdown"),
        ("beautifulsoup4 (HTML)", HAS_BS4, "pip install beautifulsoup4"),
        ("ebooklib (EPUB)", HAS_EBOOKLIB, "pip install ebooklib"),
    ]

    missing = []

    for name, available, install_cmd in deps:
        if available:
            print(f"  [OK] {name}")
        else:
            print(f"  [  ] {name} - MANCANTE")
            missing.append(install_cmd)

    # Controlla programmi esterni
    import shutil

    print("\n  Programmi esterni (opzionali):")

    externals = [
        ("LibreOffice", shutil.which("soffice") or shutil.which("libreoffice")),
        ("GIMP", shutil.which("gimp")),
        ("ImageMagick", shutil.which("convert") or shutil.which("magick")),
        ("Calibre", shutil.which("ebook-convert")),
        ("Inkscape", shutil.which("inkscape")),
    ]

    for name, path in externals:
        status = "OK" if path else "non trovato"
        symbol = "[OK]" if path else "[  ]"
        print(f"  {symbol} {name}: {status}")

    print()

    if missing:
        print("Per installare le dipendenze mancanti:")
        print(f"  pip install {' '.join(set(' '.join(missing).replace('pip install ', '').split()))}")
        print()

    return HAS_FASTAPI  # FastAPI è l'unico requisito obbligatorio


# ============================================================================
# MAIN
# ============================================================================

def main() -> None:
    """
    Entry point principale del servizio.

    Mostra banner, verifica dipendenze e avvia il server.
    """
    # Banner ASCII
    print("""
╔══════════════════════════════════════════════════════════════╗
║         DOCUMENT READER SERVICE per Open WebUI               ║
╠══════════════════════════════════════════════════════════════╣
║  Legge documenti (PDF, Word, Excel, immagini, e-book, etc.)  ║
║  e restituisce testo/JSON per uso con LLM                    ║
║                                                              ║
║  Versione: 2.0.0                                             ║
║  Autore: Carlo                                               ║
╚══════════════════════════════════════════════════════════════╝
    """)

    # Verifica dipendenze
    if not check_dependencies():
        print("[ERRORE] FastAPI è necessario per avviare il servizio.")
        print("         Installa con: pip install fastapi uvicorn python-multipart")
        sys.exit(1)

    # Info avvio
    print(f"[*] Porta: {SERVICE_PORT}")
    print(f"[*] Cache: {CACHE_DIR}")
    print(f"[*] Max file size: {MAX_FILE_SIZE_MB}MB")
    print()
    print(f"[*] Documentazione: http://localhost:{SERVICE_PORT}/docs")
    print(f"[*] Avvio servizio su http://localhost:{SERVICE_PORT}")
    print("[*] Premi Ctrl+C per fermare")
    print()

    # Avvia il server
    app = create_app()
    uvicorn.run(
        app,
        host="0.0.0.0",  # Accetta connessioni da qualsiasi IP
        port=SERVICE_PORT,
        log_level="info"
    )


# Esegui solo se chiamato direttamente (non importato)
if __name__ == "__main__":
    main()
